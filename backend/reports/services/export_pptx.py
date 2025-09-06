# =============================
# File: /app/reports/services/export_pptx.py
# Purpose: PPTX export with branded cover/ToC/content + chart slides; uses assembled["charts"].
# =============================
from __future__ import annotations
from typing import Dict, Any, List, Tuple, Optional
import os, tempfile
from datetime import datetime
from django.utils import timezone

try:
    from bs4 import BeautifulSoup
except Exception:
    BeautifulSoup = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.dml.color import RGBColor
except Exception as e:
    raise ImportError("python-pptx is required. Please install 'python-pptx'.") from e

BRAND_BLUE = RGBColor(0,123,255)
BRAND_GRAY = RGBColor(108,117,125)
DARK_GRAY  = RGBColor(73,80,87)
FONT_FAMILY = "Segoe UI"


def _add_textbox(slide, left, top, width, height, text: str, font_size_pt: int = 18, bold=False, align=PP_ALIGN.LEFT, color=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.clear(); p = tf.paragraphs[0]; run = p.add_run(); run.text = text or ""
    font = run.font; font.name = FONT_FAMILY; font.size = Pt(font_size_pt); font.bold=bool(bold)
    if color: font.color.rgb = color
    p.alignment = align
    return box


def _add_footer(slide, report_title: str, time_window: str, page_num: int, version: str = "v1"):
    t = (report_title[:40] + "...") if len(report_title) > 40 else report_title
    _add_textbox(slide, Inches(0.5), Inches(7.0), Inches(4), Inches(0.3), t, font_size_pt=10, color=BRAND_GRAY)
    if time_window:
        _add_textbox(slide, Inches(4.5), Inches(7.0), Inches(4), Inches(0.3), time_window, font_size_pt=10, align=PP_ALIGN.CENTER, color=BRAND_GRAY)
    page_info = f"Page {page_num} • {version}"
    _add_textbox(slide, Inches(8.5), Inches(7.0), Inches(3), Inches(0.3), page_info, font_size_pt=10, align=PP_ALIGN.RIGHT, color=BRAND_GRAY)


def _add_cover(prs: Presentation, title: str, time_window: str = "", generated_at: Optional[datetime]=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    if generated_at is None: generated_at = timezone.now()
    brand = slide.shapes.add_textbox(Inches(1.2), Inches(0.8), Inches(10), Inches(0.8))
    tf = brand.text_frame; p = tf.paragraphs[0]; r = p.add_run(); r.text = "MediaJira Analytics"; r.font.name = FONT_FAMILY; r.font.size = Pt(20); r.font.bold=True; r.font.color.rgb = BRAND_BLUE
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.2), Inches(1.7), Inches(11.2), Inches(1.7)); line.line.color.rgb = BRAND_BLUE; line.line.width = Pt(2)
    _add_textbox(slide, Inches(1.2), Inches(2.5), Inches(10), Inches(1.5), title, font_size_pt=44, bold=True, color=DARK_GRAY)
    subtitle = "Auto-generated Business Report" + (f"\n{time_window}" if time_window else "")
    _add_textbox(slide, Inches(1.2), Inches(4.2), Inches(10), Inches(1.0), subtitle, font_size_pt=18, color=BRAND_GRAY)
    _add_textbox(slide, Inches(1.2), Inches(6.2), Inches(10), Inches(0.4), f"Generated on {generated_at.strftime('%B %d, %Y at %I:%M %p %Z')}", font_size_pt=12, color=BRAND_GRAY)
    return slide


def _add_toc(prs: Presentation, titles: List[str], report_title: str, time_window: str, version: str):
    if not titles: return None
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_textbox(slide, Inches(1.2), Inches(0.8), Inches(10), Inches(0.8), "Table of Contents", font_size_pt=32, bold=True, color=DARK_GRAY)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.2), Inches(1.8), Inches(11.2), Inches(1.8)); line.line.color.rgb = BRAND_BLUE; line.line.width = Pt(1)
    y = 2.2
    for i, t in enumerate(titles, 1):
        _add_textbox(slide, Inches(1.5), Inches(y), Inches(10), Inches(0.4), f"{i}. {t}", font_size_pt=16, color=DARK_GRAY)
        y += 0.5
    _add_footer(slide, report_title, time_window, 2, version)
    return slide


def _add_section_slide(prs: Presentation, title: str, body_text: str, page_num: int, report_title: str, time_window: str, version: str):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_textbox(slide, Inches(1.2), Inches(0.8), Inches(10), Inches(0.8), title, font_size_pt=28, bold=True, color=DARK_GRAY)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.2), Inches(1.6), Inches(11.2), Inches(1.6)); line.line.color.rgb = BRAND_BLUE; line.line.width = Pt(1)
    _add_textbox(slide, Inches(1.2), Inches(1.9), Inches(10.5), Inches(4.5), body_text, font_size_pt=16, color=DARK_GRAY)
    _add_footer(slide, report_title, time_window, page_num, version)
    return slide


def _add_chart_slide(prs: Presentation, title: str, image_path: str, page_num: int, report_title: str, time_window: str, version: str):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_textbox(slide, Inches(1.2), Inches(0.8), Inches(10), Inches(0.6), title, font_size_pt=24, bold=True, color=DARK_GRAY)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.2), Inches(1.5), Inches(11.2), Inches(1.5)); line.line.color.rgb = BRAND_BLUE; line.line.width = Pt(1)
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1.2), Inches(1.8), width=Inches(10.0))
    else:
        _add_textbox(slide, Inches(1.2), Inches(3.0), Inches(10), Inches(2.0), "[Chart image missing]", font_size_pt=16, color=BRAND_GRAY, align=PP_ALIGN.CENTER)
    _add_footer(slide, report_title, time_window, page_num, version)
    return slide


def export_pptx(assembled: Dict[str, Any], title: str = "Report", theme: str = "light", include_raw_csv: bool = False) -> str:
    prs = Presentation()  # Use blank; for full branding, load a master: Presentation("brand_template.pptx")
    report = assembled.get("report")
    report_title = getattr(report, "title", None) or title or "Report"
    time_window = ""
    if report and getattr(report, "time_range_start", None) and getattr(report, "time_range_end", None):
        time_window = f"{report.time_range_start.strftime('%m/%d/%Y')} - {report.time_range_end.strftime('%m/%d/%Y')}"
    version = f"v{getattr(report,'id','unknown')[-8:]}" if report else "v1"

    _add_cover(prs, report_title, time_window)

    # Section titles from HTML
    html_doc: str = assembled.get("html") or ""
    section_titles: List[str] = []
    body_texts: List[str] = []
    if html_doc and BeautifulSoup is not None:
        try:
            soup = BeautifulSoup(html_doc, "html.parser")
            h2s = soup.find_all("h2")
            for h2 in h2s:
                title_txt = h2.get_text(strip=True)
                # Collect text under this h2 until next h2
                frags: List[str] = []
                for sib in h2.next_siblings:
                    if getattr(sib, "name", None) == "h2":
                        break
                    frags.append(getattr(sib, "get_text", lambda *a, **k: str(sib))("\n", strip=True))
                body_texts.append("\n".join([t for t in frags if t]))
                section_titles.append(title_txt)
        except Exception:
            section_titles = []

    has_toc = bool(section_titles)
    if has_toc:
        _add_toc(prs, section_titles, report_title, time_window, version)

    # Page numbering: 1 cover + (1 if toc)
    current_page = 1 + (1 if has_toc else 0) + 1
    for i, title_txt in enumerate(section_titles or []):
        body = body_texts[i] if i < len(body_texts) else ""
        _add_section_slide(prs, title_txt, body, current_page, report_title, time_window, version)
        current_page += 1

    # Chart slides from assembled metadata
    for i, c in enumerate(assembled.get("charts") or []):
        _add_chart_slide(prs, c.get("title") or f"Chart {i+1}", c.get("path"), current_page, report_title, time_window, version)
        current_page += 1

    # CSV appendix (optional) – you can wire in assembled["csv_paths"] if you generate them
    # (omitted here for brevity – same approach as earlier)

    fd, out_path = tempfile.mkstemp(suffix=".pptx"); os.close(fd)
    try:
        prs.save(out_path)
    except Exception:
        os.unlink(out_path); raise
    return out_path


